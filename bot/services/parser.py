"""Fayl parserlari - DOC, DOCX, PDF, XLSX, TXT."""
import io
import re
import logging
from pathlib import Path
from typing import Optional
import chardet

# PDF
import pdfplumber

# DOCX
from docx import Document

# XLSX
from openpyxl import load_workbook

# DOC (eski format) - antiword yordamida (serverda o'rnatilgan bo'lishi kerak)
# Eslatma: antiword bo'lmasa, DOC fayllar o'qilmaydi. Tavsiya: DOCX ga o'tkazish.
import subprocess

logger = logging.getLogger(__name__)


class FileParser:
    """Universal fayl parser - barcha formatlarni qo'llab-quvvatlaydi."""

    SUPPORTED = {"pdf", "doc", "docx", "xlsx", "xls", "txt", "pptx"}

    # O'rtacha tokenga yaqin matn uzunligi: 1 token ~= 4 belgi
    MAX_CHARS = 60_000  # 60K belgi - katta fayllar uchun ham yetarli

    @classmethod
    async def extract_text(cls, file_path: str | Path, file_ext: str) -> str:
        """Fayldan matn ajratib olish. Asosiy entry point."""
        file_ext = file_ext.lower().lstrip(".")
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"Fayl topilmadi: {file_path}")

        # Formatga qarab mos parseerga yo'naltirish
        extractor_map = {
            "pdf": cls._extract_pdf,
            "docx": cls._extract_docx,
            "doc": cls._extract_doc,
            "xlsx": cls._extract_xlsx,
            "xls": cls._extract_xls,
            "txt": cls._extract_txt,
            "pptx": cls._extract_pptx,
        }

        extractor = extractor_map.get(file_ext)
        if not extractor:
            raise ValueError(f"Qo'llab-quvvatlanmaydigan format: {file_ext}")

        text = await extractor(path)

        # Matnni tozalash
        text = cls._clean_text(text)
        return text[: cls.MAX_CHARS]

    @staticmethod
    async def _extract_pdf(path: Path) -> str:
        """PDF dan matn ajratish. pdfplumber yordamida (jadvallar bilan yaxshi ishlaydi)."""
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        text_parts.append(f"\n[--- Sahifa {page_num} ---]\n{page_text}")

                    # Jadvallar
                    tables = page.extract_tables()
                    for table_idx, table in enumerate(tables, 1):
                        if table:
                            text_parts.append(f"\n[Jadval {page_num}.{table_idx}]")
                            for row in table:
                                if row:
                                    text_parts.append(" | ".join(str(c) for c in row if c))
                except Exception as e:
                    logger.warning(f"Sahifa {page_num} o'qib bo'lmadi: {e}")
                    continue

        return "\n".join(text_parts)

    @staticmethod
    async def _extract_docx(path: Path) -> str:
        """DOCX dan matn (rasmlar, jadvallar, header/footer bilan)."""
        doc = Document(path)
        text_parts = []

        # Header / Footer
        for section in doc.sections:
            if section.header:
                header_text = "\n".join(p.text for p in section.header.paragraphs if p.text.strip())
                if header_text:
                    text_parts.append(f"[HEADER]\n{header_text}\n[/HEADER]")
            if section.footer:
                footer_text = "\n".join(p.text for p in section.footer.paragraphs if p.text.strip())
                if footer_text:
                    text_parts.append(f"[FOOTER]\n{footer_text}\n[/FOOTER]")

        # Asosiy kontent
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Jadvallar
        for table_idx, table in enumerate(doc.tables, 1):
            text_parts.append(f"\n[Jadval {table_idx}]")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    text_parts.append(" | ".join(cells))

        return "\n".join(text_parts)

    @staticmethod
    async def _extract_doc(path: Path) -> str:
        """Eski DOC format - antiword yoki catdoc yordamida (serverda o'rnatilgan bo'lishi kerak)."""
        try:
            result = subprocess.run(
                ["antiword", str(path)],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout:
                return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Fallback - mammoth
        try:
            import mammoth
            with open(path, "rb") as f:
                result = mammoth.extract_raw_text(f)
                return result.value
        except Exception as e:
            logger.error(f"DOC o'qib bo'lmadi: {e}")
            raise RuntimeError("DOC formatini o'qib bo'lmadi. Iltimos DOCX yoki PDF ishlating.")

    @staticmethod
    async def _extract_xlsx(path: Path) -> str:
        """XLSX dan barcha varaqlar."""
        wb = load_workbook(path, data_only=True, read_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n[=== Varaq: {sheet_name} ===]")

            for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
                if any(cell is not None and str(cell).strip() for cell in row):
                    cells = [str(c).strip() if c is not None else "" for c in row]
                    text_parts.append(" | ".join(c for c in cells if c))

                    if row_idx > 5000:  # Katta fayllar uchun limit
                        text_parts.append("... (qolgan qatorlar qisqartirildi)")
                        break

        wb.close()
        return "\n".join(text_parts)

    @staticmethod
    async def _extract_xls(path: Path) -> str:
        """Eski XLS - xlrd yoki fallback to libreoffice."""
        try:
            import xlrd
            wb = xlrd.open_workbook(str(path))
            text_parts = []
            for sheet in wb.sheets():
                text_parts.append(f"\n[=== Varaq: {sheet.name} ===]")
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    cells = [str(c).strip() for c in row if str(c).strip()]
                    if cells:
                        text_parts.append(" | ".join(cells))
            return "\n".join(text_parts)
        except ImportError:
            raise RuntimeError("XLS formatini o'qish uchun xlrd kerak")

    @staticmethod
    async def _extract_txt(path: Path) -> str:
        """TXT - encoding detect bilan."""
        with open(path, "rb") as f:
            raw = f.read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "utf-8"

        try:
            return raw.decode(encoding, errors="ignore")
        except (LookupError, UnicodeDecodeError):
            return raw.decode("utf-8", errors="ignore")

    @staticmethod
    async def _extract_pptx(path: Path) -> str:
        """PPTX (PowerPoint) slaydlardan matn."""
        try:
            from pptx import Presentation
            pres = Presentation(path)
            text_parts = []
            for slide_num, slide in enumerate(pres.slides, 1):
                text_parts.append(f"\n[--- Slayd {slide_num} ---]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except ImportError:
            raise RuntimeError("PPTX uchun python-pptx kerak")

    @staticmethod
    def _clean_text(text: str) -> str:
        """Matnni AI uchun tayyorlash - keraksiz bo'shliqlar, belgilarni tozalash."""
        if not text:
            return ""

        # Null bytes va boshqa yashirin belgilar
        text = text.replace("\x00", "").replace("\r", "\n")

        # Ortiqcha bo'sh qatorlar (3+ ketma-ket \n ni 2 taga)
        text = re.sub(r"\n{3,}", "\n\n", text)

        # Tab va bo'shliqlarni birlashtirish
        text = re.sub(r"[ \t]+", " ", text)

        # Qator boshida va oxiridagi bo'shliqlar
        text = "\n".join(line.strip() for line in text.split("\n"))

        # Juda qisqa qatorlarni (1-2 belgi) olib tashlash
        lines = [l for l in text.split("\n") if len(l) > 2 or not l.strip()]
        text = "\n".join(lines)

        return text.strip()

    @classmethod
    def detect_format(cls, filename: str) -> Optional[str]:
        """Fayl nomidan format aniqlash."""
        ext = Path(filename).suffix.lower().lstrip(".")
        return ext if ext in cls.SUPPORTED else None

    @classmethod
    def get_file_stats(cls, text: str) -> dict:
        """Matn statistikasi (analytics uchun)."""
        if not text:
            return {"chars": 0, "words": 0, "lines": 0, "sentences": 0}

        # O'zbek/Rus/Ingliz tillar uchun so'zlar
        words = re.findall(r"\b\w+\b", text, re.UNICODE)

        # Gaplar
        sentences = re.split(r"[.!?]+\s+", text)

        return {
            "chars": len(text),
            "words": len(words),
            "lines": text.count("\n") + 1,
            "sentences": len([s for s in sentences if s.strip()]),
        }
